"""
ModelGate — PPTX Builder (v2)
Clean, professional dark-themed PowerPoint.
Uses tables for data layouts instead of manual positioning.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn

# ── Palette ────────────────────────────────────────────────────────
BG      = RGBColor(0x10, 0x10, 0x10)
CARD    = RGBColor(0x1a, 0x1a, 0x1a)
CARD2   = RGBColor(0x22, 0x22, 0x22)
BORDER  = RGBColor(0x33, 0x33, 0x33)
TEXT    = RGBColor(0xee, 0xee, 0xee)
TEXT2   = RGBColor(0x99, 0x99, 0x99)
TEXT3   = RGBColor(0x5a, 0x5a, 0x5a)
GHOST   = RGBColor(0x20, 0x20, 0x20)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

SRF = 'Georgia'
SAN = 'Calibri'
MON = 'Consolas'

W, H = 13.333, 7.5
ML, MR, MT, MB = 0.8, 0.8, 0.6, 0.6


# ══════════════════════════════════════════════════════════════════
#  HELPERS
# ══════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    return prs

def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def text(slide, txt, left, top, width, height,
         font=SAN, size=14, bold=False, italic=False,
         color=TEXT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tf

def add_p(tf, txt, font=SAN, size=14, bold=False, italic=False,
          color=TEXT, align=PP_ALIGN.LEFT, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = txt
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p

def box(slide, left, top, width, height, fill_color=CARD, border_color=BORDER):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def line(slide, left, top, width, color=BORDER):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(left), Inches(top), Inches(left + width), Inches(top))
    c.line.color.rgb = color
    c.line.width = Pt(0.5)

def make_table(slide, data, left, top, width, col_widths,
               header=True, row_height=0.45):
    rows = len(data)
    cols = len(col_widths)
    tbl_shape = slide.shapes.add_table(rows, cols,
        Inches(left), Inches(top), Inches(width), Inches(rows * row_height))
    tbl = tbl_shape.table

    # Set column widths
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)

    # Style cells
    for r_idx, row_data in enumerate(data):
        is_hdr = header and r_idx == 0
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = ""

            # Cell background
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn('a:solidFill'), {})
            srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': '101010' if not is_hdr else '1a1a1a'})
            solidFill.append(srgb)
            tcPr.append(solidFill)

            # Borders
            for edge in ['lnL', 'lnR', 'lnT', 'lnB']:
                ln = tcPr.makeelement(qn('a:' + edge), {'w': '6350'})
                sf = ln.makeelement(qn('a:solidFill'), {})
                sc = sf.makeelement(qn('a:srgbClr'), {'val': '333333'})
                sf.append(sc)
                ln.append(sf)
                tcPr.append(ln)

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(cell_text)

            if is_hdr:
                r.font.name = MON
                r.font.size = Pt(9)
                r.font.bold = False
                r.font.color.rgb = TEXT3
            else:
                r.font.name = SAN
                r.font.size = Pt(12)
                r.font.bold = (c_idx == 0)
                r.font.color.rgb = TEXT if c_idx == 0 else TEXT2

            cell.text_frame.margin_left = Inches(0.08)
            cell.text_frame.margin_right = Inches(0.08)
            cell.text_frame.margin_top = Inches(0.06)
            cell.text_frame.margin_bottom = Inches(0.06)
            cell.vertical_anchor = MSO_ANCHOR.TOP

    return tbl_shape

def slide_header(slide, num, eyebrow_text, title_text, subtitle=None):
    """Standard header: eyebrow + title + optional subtitle."""
    # Slide number
    text(slide, f'{num:02d}', ML, MT - 0.1, 1.5, 0.5,
         font=MON, size=11, color=TEXT3)
    # Eyebrow
    text(slide, eyebrow_text.upper(), ML, MT + 0.35, 8, 0.25,
         font=MON, size=9, color=TEXT3)
    # Title
    text(slide, title_text, ML, MT + 0.7, W - ML - MR, 0.8,
         font=SRF, size=30, color=TEXT)
    if subtitle:
        text(slide, subtitle, ML, MT + 1.2, W - ML - MR, 0.4,
             font=SAN, size=13, italic=True, color=TEXT3)
    line(slide, ML, MT + (1.7 if subtitle else 1.35), W - ML - MR, color=BORDER)
    content_top = MT + (1.85 if subtitle else 1.5)
    return content_top


# ══════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════

def s01_title(prs):
    slide = add_slide(prs)

    # Center area
    text(slide, 'KSU SOCIAL GOOD HACKATHON 2026  ·  ASSURANT TRACK',
         ML, 1.5, W - ML - MR, 0.3,
         font=MON, size=10, color=TEXT3, align=PP_ALIGN.CENTER)

    text(slide, 'ModelGate', ML, 2.1, W - ML - MR, 1.5,
         font=SRF, size=80, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    line(slide, W/2 - 0.5, 3.4, 1.0, color=BORDER)

    text(slide, 'Contract-Aware AI Control Plane',
         ML, 3.6, W - ML - MR, 0.5,
         font=SAN, size=18, color=TEXT2, align=PP_ALIGN.CENTER)

    lines = ['One line of code.  Immediate savings.  No friction.']
    text(slide, lines[0], ML, 4.3, W - ML - MR, 0.4,
         font=SAN, size=14, color=TEXT3, align=PP_ALIGN.CENTER)

    # Tech tags along bottom
    tags = ['FastAPI + Python', 'Next.js + TypeScript',
            'Arch-Router-1.5B', 'OpenRouter', 'Docker']
    tag_text = '   ·   '.join(tags)
    text(slide, tag_text, ML, H - MB - 0.4, W - ML - MR, 0.3,
         font=MON, size=9, color=TEXT3, align=PP_ALIGN.CENTER)


def s02_assurant(prs):
    slide = add_slide(prs)
    ct = slide_header(slide, 2, 'The Assurant Way',
                      'Helping people thrive in a connected world.',
                      'The connected world runs on AI. AI is expensive, wasteful, '
                      'and impossible to configure correctly at scale.')

    # 2x2 stat grid
    stats = [
        ('180×',   'more energy consumed by premium models\nvs small models per query'),
        ('50–90%', 'of enterprise AI inference spend wasted\non overprovisioned models'),
        ('$37B',   'enterprise GenAI spend in 2025\nup 3.2× year-over-year'),
        ('51%',    'of organizations can actually\nmeasure their AI ROI'),
    ]
    gw = (W - ML - MR - 0.25) / 2
    gh = (H - MB - ct - 0.65) / 2

    for i, (num, desc) in enumerate(stats):
        c, r = i % 2, i // 2
        gl = ML + c * (gw + 0.25)
        gt = ct + r * (gh + 0.15)
        box(slide, gl, gt, gw, gh, fill_color=CARD, border_color=BORDER)
        text(slide, num, gl + 0.3, gt + 0.25, gw - 0.6, 0.8,
             font=SRF, size=44, bold=True, color=TEXT)
        text(slide, desc, gl + 0.3, gt + 1.05, gw - 0.6, 0.8,
             font=SAN, size=12, color=TEXT2)

    text(slide, '"Without security and efficient application infrastructure, we are at risk '
         'of increasing the waste and consumption with our day-to-day tools." — Assurant Challenge',
         ML, H - MB - 0.35, W - ML - MR, 0.3,
         font=SAN, size=10, italic=True, color=TEXT3)


def s03_problem(prs):
    slide = add_slide(prs)
    ct = slide_header(slide, 3, 'The Problem Is Universal',
                      'Every layer of the AI stack is affected.')

    data = [
        ['AFFECTED PARTY', 'THE PROBLEM TODAY', 'CONSEQUENCE'],
        ['Developers &\nProduct Teams',
         'No time to evaluate 200+ models. Pick one premium\noption and never revisit it.',
         'Overpay by\n10–30×'],
        ['Enterprise\nCustomers',
         'Contract constraints manually translated into config\n— or silently ignored. Compliance risk.',
         'Legal\nexposure'],
        ['End Users',
         '"What is your return policy?" routed to a slow premium\nmodel. Waits 2 seconds for a simple answer.',
         'Poor\nexperience'],
        ['AI Data\nCenters',
         'Unnecessary load on premium GPU clusters from queries\nthat a small model would answer perfectly.',
         'Wasted\ncapacity'],
        ['The\nEnvironment',
         'Premium models use ~39 Wh/query vs ~0.22 Wh for small\nmodels. A 180× energy difference per call.',
         '415 TWh/yr\nand rising'],
    ]

    make_table(slide, data, ML, ct, W - ML - MR,
               [2.0, 7.2, 1.9], row_height=0.7)


def s04_solution(prs):
    slide = add_slide(prs)
    ct = slide_header(slide, 4, 'The Fix',
                      'One line of code.  We handle everything else.')

    # Two code boxes side by side
    cw = (W - ML - MR - 0.4) / 2

    # BEFORE box
    box(slide, ML, ct, cw, 2.55, fill_color=CARD, border_color=BORDER)
    box(slide, ML, ct, cw, 0.35, fill_color=CARD2, border_color=BORDER)
    text(slide, 'BEFORE — EVERY TEAM TODAY', ML + 0.15, ct + 0.07, cw - 0.3, 0.25,
         font=MON, size=9, color=TEXT3)

    before_code = (
        "# Direct provider call\n"
        "from openai import OpenAI\n"
        "\n"
        "client = OpenAI(\n"
        '    api_key="sk-...",\n'
        '    base_url="https://api.openai.com/v1"\n'
        ")\n"
        "\n"
        "# Same model. Every request. Forever."
    )
    text(slide, before_code, ML + 0.2, ct + 0.45, cw - 0.4, 2.0,
         font=MON, size=11, color=TEXT2)

    # AFTER box
    al = ML + cw + 0.4
    box(slide, al, ct, cw, 2.55, fill_color=CARD2, border_color=RGBColor(0x44, 0x44, 0x44))
    box(slide, al, ct, cw, 0.35, fill_color=RGBColor(0x2a, 0x2a, 0x2a),
        border_color=RGBColor(0x44, 0x44, 0x44))
    text(slide, 'AFTER — WITH MODELGATE', al + 0.15, ct + 0.07, cw - 0.3, 0.25,
         font=MON, size=9, color=TEXT2)

    after_code = (
        "# Same code. Different URL.\n"
        "from openai import OpenAI\n"
        "\n"
        "client = OpenAI(\n"
        '    api_key="<token>",\n'
        '    base_url="https://gate/acme/v1"  # <--\n'
        ")\n"
        "\n"
        "# Right model. Every request. Auto."
    )
    text(slide, after_code, al + 0.2, ct + 0.45, cw - 0.4, 2.0,
         font=MON, size=11, color=TEXT)

    # Bottom: 3 benefits
    bt = ct + 2.8
    bw = (W - ML - MR - 0.5) / 3
    benefits = [
        ('API COMPATIBLE', 'Full OpenAI spec — no SDK\nchanges, no rewrites.'),
        ('POLICY ENFORCED', 'Contract constraints become\nrouting rules, automatically.'),
        ('ALWAYS OPTIMAL', 'Cheapest model that satisfies\nevery constraint, per request.'),
    ]
    for i, (label, desc) in enumerate(benefits):
        bl = ML + i * (bw + 0.25)
        box(slide, bl, bt, bw, 1.05, fill_color=CARD, border_color=BORDER)
        text(slide, label, bl + 0.2, bt + 0.15, bw - 0.4, 0.25,
             font=MON, size=9, bold=True, color=TEXT2)
        text(slide, desc, bl + 0.2, bt + 0.5, bw - 0.4, 0.5,
             font=SAN, size=12, color=TEXT2)


def s05_how(prs):
    slide = add_slide(prs)
    ct = slide_header(slide, 5, 'System Architecture',
                      'Two phases. Fully automated. No ongoing effort.')

    cw = (W - ML - MR - 0.5) / 2

    # ── Phase 1 ──
    box(slide, ML, ct, cw, 0.35, fill_color=CARD2, border_color=BORDER)
    text(slide, 'PHASE 1:  ONBOARDING  (≈ 30 seconds)',
         ML + 0.2, ct + 0.07, cw - 0.4, 0.25, font=MON, size=9, color=TEXT2)

    steps = [
        'Upload contract  (PDF / text / SLA docs)',
        'LLM reads and extracts all constraints',
        'Structured Customer AI Profile generated',
        'OpenAI-compatible endpoint goes live',
    ]
    for i, step in enumerate(steps):
        st = ct + 0.5 + i * 0.7
        box(slide, ML, st, cw, 0.55, fill_color=CARD, border_color=BORDER)
        text(slide, f'{i+1}', ML + 0.15, st + 0.08, 0.3, 0.35,
             font=MON, size=16, bold=True, color=TEXT3)
        text(slide, step, ML + 0.5, st + 0.12, cw - 0.7, 0.35,
             font=SAN, size=13, color=TEXT2)

    # ── Phase 2 ──
    l2 = ML + cw + 0.5
    box(slide, l2, ct, cw, 0.35, fill_color=CARD2, border_color=BORDER)
    text(slide, 'PHASE 2:  RUNTIME ROUTING  (~50ms overhead)',
         l2 + 0.2, ct + 0.07, cw - 0.4, 0.25, font=MON, size=9, color=TEXT2)

    # Pipeline arrow
    pipeline = ['Prompt', 'Classify', 'Filter', 'Score', 'Route', 'Respond']
    pw = (cw - 0.3) / len(pipeline)
    for i, step in enumerate(pipeline):
        pl = l2 + 0.15 + i * pw
        box(slide, pl, ct + 0.5, pw - 0.06, 0.32,
            fill_color=CARD2 if i in (0, 5) else CARD,
            border_color=BORDER)
        text(slide, step, pl + 0.04, ct + 0.55, pw - 0.14, 0.22,
             font=MON, size=8, color=TEXT if i in (0, 5) else TEXT2,
             align=PP_ALIGN.CENTER)

    # Tier table
    tier_data = [
        ['TIER', 'MODELS', 'USE CASE'],
        ['Simple', 'GPT 5.4 nano · Gemini 3.1 Flash Lite · Haiku 4.5',
         'FAQ, greetings, status checks · ~300ms'],
        ['Medium', 'GPT 5.4 mini · Sonnet 4.6 · Grok 4.1 Fast',
         'Summarize, explain, troubleshoot · ~700ms'],
        ['Complex', 'GPT 5.4 · Opus 4.6 · Gemini 3.1 Pro · Grok 4.20',
         'Multi-doc reasoning, code gen, legal · ~1500ms'],
    ]
    make_table(slide, tier_data, l2, ct + 1.05, cw,
               [0.8, 2.8, 2.3], row_height=0.55)

    text(slide, 'Classifier: katanemo/Arch-Router-1.5B on GPU  ·  Heuristic fallback if no GPU',
         l2, ct + 3.35, cw, 0.3, font=MON, size=9, color=TEXT3)


def s06_personas(prs):
    slide = add_slide(prs)
    ct = slide_header(slide, 6, 'Who It Serves',
                      'Three personas. One solution. Simultaneous benefit.')

    cw = (W - ML - MR - 0.5) / 3
    ch = H - MB - ct - 0.05

    personas = [
        ('01', 'DEVELOPER / PRODUCT TEAM',
         '"I just want AI to work.\nI don\'t have time to evaluate\nevery new model."',
         'Picks one premium model,\noverpays forever.',
         'One URL change. Automatic\nright-sizing. Done forever.'),
        ('02', 'ENTERPRISE COMPLIANCE',
         '"Our contracts say EU-only.\nI cannot trust engineers\nwill remember this."',
         'Manual config per customer.\nOne mistake = violation.',
         'Upload the contract. Routing\nenforces it automatically.'),
        ('03', 'OPERATIONS / PLATFORM',
         '"I have no idea what our AI\nis costing us or which\nmodels we\'re using."',
         'No visibility, no cost tracking,\nno model distribution data.',
         'Live dashboard. Cost savings.\nPer-request explanations.'),
    ]

    for i, (num, role, quote, before, after) in enumerate(personas):
        cl = ML + i * (cw + 0.25)

        # Card background
        box(slide, cl, ct, cw, ch, fill_color=CARD, border_color=BORDER)

        # Number
        text(slide, num, cl + 0.25, ct + 0.15, 1.0, 0.6,
             font=SRF, size=48, bold=True, color=GHOST)

        # Role
        text(slide, role, cl + 0.25, ct + 0.7, cw - 0.5, 0.25,
             font=MON, size=9, color=TEXT3)

        line(slide, cl + 0.15, ct + 1.05, cw - 0.3, color=BORDER)

        # Quote
        text(slide, quote, cl + 0.25, ct + 1.2, cw - 0.5, 1.1,
             font=SRF, size=13, italic=True, color=TEXT2)

        line(slide, cl + 0.15, ct + 2.4, cw - 0.3, color=BORDER)

        # Before
        text(slide, 'BEFORE', cl + 0.25, ct + 2.55, 1.0, 0.22,
             font=MON, size=8, color=TEXT3)
        text(slide, before, cl + 0.25, ct + 2.8, cw - 0.5, 0.8,
             font=SAN, size=11, color=TEXT2)

        line(slide, cl + 0.15, ct + 3.65, cw - 0.3, color=BORDER)

        # After
        text(slide, 'AFTER', cl + 0.25, ct + 3.8, 1.0, 0.22,
             font=MON, size=8, color=TEXT2)
        text(slide, after, cl + 0.25, ct + 4.05, cw - 0.5, 0.8,
             font=SAN, size=11, color=TEXT)


def s07_evidence(prs):
    slide = add_slide(prs)
    ct = slide_header(slide, 7, 'Evidence',
                      'Academic validation + live system results.')

    # ── Left: Research ──
    lw = (W - ML - MR) * 0.55
    text(slide, 'PUBLISHED RESEARCH', ML, ct, lw, 0.25,
         font=MON, size=9, color=TEXT3)

    research_data = [
        ['SOURCE', 'FINDING'],
        ['FrugalGPT\n(Stanford, 2023)',
         'Up to 98% cost reduction while matching\nGPT-4 quality by cascading through cheaper models.'],
        ['RouteLLM\n(LMSYS, ICLR 2025)',
         '>85% cost reduction at 95% of GPT-4\nperformance using a trained routing model.'],
        ['RouteLLM Matrix\nFactorization',
         'Only 14% of queries actually require the most\ncapable model. 86% can use smaller ones.'],
        ['IEA Energy\nReport (2025)',
         'Global data center consumption: 415 TWh in 2024.\nProjected to exceed 1,000 TWh by 2026.'],
    ]
    make_table(slide, research_data, ML, ct + 0.3, lw,
               [1.8, lw - 1.8], row_height=0.65)

    # ── Right: Live metrics ──
    rl = ML + lw + 0.4
    rw = W - MR - rl
    text(slide, 'MODELGATE LIVE RESULTS', rl, ct, rw, 0.25,
         font=MON, size=9, color=TEXT3)

    metrics = [
        ('Onboarding', '< 30s', 'contract → endpoint'),
        ('Classify', '~50ms', 'Arch-Router-1.5B'),
        ('Simple %', '  %', 'fill after demo'),
        ('Savings', '  %', 'vs always-premium'),
    ]

    mw = (rw - 0.15) / 2
    mh = 1.3
    for i, (label, value, note) in enumerate(metrics):
        mc, mr = i % 2, i // 2
        ml = rl + mc * (mw + 0.15)
        mt = ct + 0.35 + mr * (mh + 0.15)
        box(slide, ml, mt, mw, mh, fill_color=CARD, border_color=BORDER)
        text(slide, label.upper(), ml + 0.18, mt + 0.12, mw - 0.36, 0.2,
             font=MON, size=8, color=TEXT3)
        text(slide, value, ml + 0.18, mt + 0.4, mw - 0.36, 0.6,
             font=SRF, size=36, bold=True, color=TEXT)
        text(slide, note, ml + 0.18, mt + 1.0, mw - 0.36, 0.2,
             font=MON, size=8, color=TEXT3)

    # Quote box
    box(slide, rl, ct + 3.05, rw, 0.75, fill_color=CARD, border_color=BORDER)
    text(slide, '"Smart routing delivers 85–98% cost reduction while maintaining '
         'quality indistinguishable from always-premium."',
         rl + 0.2, ct + 3.15, rw - 0.4, 0.6,
         font=SRF, size=12, italic=True, color=TEXT2)


def s08_close(prs):
    slide = add_slide(prs)
    ct = slide_header(slide, 8, 'Impact & Close',
                      'One change. Benefits at every layer of the stack.')

    impact_data = [
        ['LAYER', 'HOW MODELGATE HELPS', 'HEADLINE'],
        ['Business', '60–98% cost savings on AI inference\nthrough intelligent right-sizing', '60–98% saved'],
        ['End Users', 'Simple queries get fast answers.\nComplex queries get powerful models.', 'Right model, always'],
        ['Data Centers', 'Less unnecessary GPU load on premium\nclusters. Resources freed.', 'Less waste'],
        ['Energy Grid', 'Up to 180× less energy per routed-away\npremium call. Compounding effect.', '180× difference'],
        ['Compliance', 'Contract constraints become routing rules.\nEnforced per request, automatically.', 'Zero human error'],
    ]
    make_table(slide, impact_data, ML, ct, W - ML - MR,
               [1.6, 6.4, 2.1], row_height=0.6)

    # Closing
    bt = ct + 3.8
    line(slide, ML, bt, W - ML - MR, color=BORDER)

    text(slide, '"ModelGate is the lowest-friction AI optimization available today. '
         'You change one line of code. We eliminate thousands of wasteful premium model calls — '
         'automatically."',
         ML, bt + 0.15, (W - ML - MR) * 0.6, 0.8,
         font=SRF, size=14, italic=True, color=TEXT2)

    text(slide, 'Thank you.', ML + (W - ML - MR) * 0.6 + 0.3, bt + 0.1,
         (W - ML - MR) * 0.38, 0.7,
         font=SRF, size=44, bold=True, color=TEXT, align=PP_ALIGN.RIGHT)

    text(slide, 'KSU  ·  Assurant Track  ·  2026',
         ML + (W - ML - MR) * 0.6 + 0.3, bt + 0.7,
         (W - ML - MR) * 0.38, 0.3,
         font=MON, size=9, color=TEXT3, align=PP_ALIGN.RIGHT)


def s09_demo(prs):
    slide = add_slide(prs)
    ct = slide_header(slide, 9, 'Live Demo',
                      'Five steps. One contract. Immediate routing.')

    steps = [
        ('1', 'Onboard ACME\nSupport',
         'Upload contract\nExtract Profile',
         'EU-only · DeepSeek blocked\n1000ms latency target\nHigh privacy · Auto tiers'),
        ('2', 'Copy\nEndpoint',
         'Profile page\nCopy URL',
         '/acme-support/v1\n/chat/completions\nDrop-in for OpenAI'),
        ('3', 'Simple\nPrompt',
         '"What is your\nreturn policy?"',
         'Classified: simple\nGPT 5.4 nano · ~300ms\nFraction of a cent'),
        ('4', 'Complex\nPrompt',
         '"Analyze liability\nexposure..."',
         'Classified: complex\nOpus 4.6 · ~1.5s\nSame endpoint, auto'),
        ('5', 'Dashboard\nMetrics',
         'Navigate to\nmain dashboard',
         'Model distribution chart\nCost savings banner\nLive routing feed'),
    ]

    cw = (W - ML - MR - 0.4) / 5
    ch = H - MB - ct - 0.1

    for i, (num, title_t, action, expects) in enumerate(steps):
        cl = ML + i * (cw + 0.1)
        box(slide, cl, ct, cw, ch, fill_color=CARD, border_color=BORDER)

        text(slide, num, cl + 0.2, ct + 0.1, 0.5, 0.5,
             font=SRF, size=36, bold=True, color=GHOST)
        text(slide, title_t, cl + 0.2, ct + 0.6, cw - 0.4, 0.65,
             font=SAN, size=13, bold=True, color=TEXT)

        line(slide, cl + 0.1, ct + 1.35, cw - 0.2, color=BORDER)

        text(slide, 'ACTION', cl + 0.2, ct + 1.5, cw - 0.4, 0.2,
             font=MON, size=8, color=TEXT3)
        text(slide, action, cl + 0.2, ct + 1.72, cw - 0.4, 0.7,
             font=MON, size=10, color=TEXT2)

        line(slide, cl + 0.1, ct + 2.55, cw - 0.2, color=BORDER)

        text(slide, 'EXPECTED', cl + 0.2, ct + 2.7, cw - 0.4, 0.2,
             font=MON, size=8, color=TEXT3)
        text(slide, expects, cl + 0.2, ct + 2.92, cw - 0.4, 1.5,
             font=SAN, size=10, color=TEXT2)


def s10_planning(prs):
    slide = add_slide(prs)
    ct = slide_header(slide, 10, 'Planning & Execution',
                      'MVP scoped, artifacts produced, team distributed.')

    cw = (W - ML - MR - 0.4) / 2

    # ── Left: MVP delivered ──
    text(slide, 'MVP DELIVERED', ML, ct, cw, 0.22, font=MON, size=9, color=TEXT3)
    line(slide, ML, ct + 0.28, cw, color=BORDER)

    mvp = [
        'Contract ingestion + LLM extraction pipeline',
        'Structured CustomerProfile JSON schema',
        'OpenAI-compatible proxy endpoint',
        'Arch-Router-1.5B classifier + heuristic fallback',
        'Policy-filtered, objective-scored routing engine',
        'Dashboard: stats, profiles, logs, playground',
    ]
    tf = text(slide, '', ML, ct + 0.38, cw, 2.5, font=SAN, size=12, color=TEXT2)
    for i, item in enumerate(mvp):
        if i == 0:
            tf.paragraphs[0].add_run().text = f'  {item}'
            tf.paragraphs[0].runs[0].font.name = SAN
            tf.paragraphs[0].runs[0].font.size = Pt(12)
            tf.paragraphs[0].runs[0].font.color.rgb = TEXT2
        else:
            add_p(tf, f'  {item}', font=SAN, size=12, color=TEXT2, space_before=6)

    # Next iterations
    ni_t = ct + 2.95
    text(slide, 'NEXT ITERATIONS', ML, ni_t, cw, 0.22, font=MON, size=9, color=TEXT3)
    line(slide, ML, ni_t + 0.28, cw, color=BORDER)

    nxt = [
        'Vector search over large contract doc sets',
        'Live provider health signals in routing',
        'Managed SaaS deployment (no Docker)',
        'Streaming response support',
        'Auto model catalog updates',
    ]
    tf2 = text(slide, '', ML, ni_t + 0.38, cw, 2.0, font=SAN, size=11, color=TEXT3)
    for i, item in enumerate(nxt):
        if i == 0:
            tf2.paragraphs[0].add_run().text = f'  {item}'
            tf2.paragraphs[0].runs[0].font.name = SAN
            tf2.paragraphs[0].runs[0].font.size = Pt(11)
            tf2.paragraphs[0].runs[0].font.color.rgb = TEXT3
        else:
            add_p(tf2, f'  {item}', font=SAN, size=11, color=TEXT3, space_before=5)

    # ── Right: Artifacts + Team ──
    l2 = ML + cw + 0.4

    text(slide, 'DESIGN ARTIFACTS (5)', l2, ct, cw, 0.22, font=MON, size=9, color=TEXT3)
    line(slide, l2, ct + 0.28, cw, color=BORDER)

    artifacts_data = [
        ['#', 'ARTIFACT', 'DESCRIPTION'],
        ['01', 'Architecture Diagram', 'Backend services + routing pipeline'],
        ['02', 'Profile Schema', 'JSON: constraints, performance, routing'],
        ['03', 'Routing Flowchart', 'Policy filter → score → select'],
        ['04', 'Dashboard UI', 'Overview, customers, profile, logs, playground'],
        ['05', 'Demo Script', 'Step-by-step with expected outputs'],
    ]
    make_table(slide, artifacts_data, l2, ct + 0.38, cw,
               [0.35, 1.8, cw - 2.15], row_height=0.4)

    # Team
    tm_t = ct + 2.95
    text(slide, 'TEAM DISTRIBUTION (4 MEMBERS)', l2, tm_t, cw, 0.22,
         font=MON, size=9, color=TEXT3)
    line(slide, l2, tm_t + 0.28, cw, color=BORDER)

    team_data = [
        ['DOMAIN', 'SCOPE'],
        ['Backend Routing', 'Classifier, router engine, provider registry'],
        ['LLM Extraction', 'Extractor, prompt eng, profile schema'],
        ['Frontend', 'Next.js pages, shadcn/ui, Recharts'],
        ['Integration', 'Docker, OpenRouter, e2e tests, demo prep'],
    ]
    make_table(slide, team_data, l2, tm_t + 0.38, cw,
               [1.6, cw - 1.6], row_height=0.4)


# ══════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    builders = [
        s01_title, s02_assurant, s03_problem, s04_solution,
        s05_how, s06_personas, s07_evidence, s08_close,
        s09_demo, s10_planning,
    ]
    for b in builders:
        b(prs)
        print(f'  built {b.__name__}')

    out = '/home/cbak/programming/ksu-sg-hackathon/docs/modelgate-presentation.pptx'
    prs.save(out)
    print(f'\nSaved → {out}')

if __name__ == '__main__':
    main()
